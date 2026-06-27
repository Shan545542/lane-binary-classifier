from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "智能驾驶车道二分类PPT设计文档.pptx"

WIDE_WIDTH = Inches(13.333)
WIDE_HEIGHT = Inches(7.5)

BG = RGBColor(248, 250, 252)
INK = RGBColor(20, 31, 43)
MUTED = RGBColor(83, 96, 112)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(234, 88, 12)
LINE = RGBColor(203, 213, 225)
WHITE = RGBColor(255, 255, 255)


def set_run(run, size: int = 18, bold: bool = False, color: RGBColor = INK) -> None:
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.55))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run(run, 24, True)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.93), Inches(11.8), Inches(0.35))
        p2 = sub.text_frame.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        set_run(run2, 11, False, MUTED)


def add_footer(slide, page: int) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(7.12), Inches(12.25), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    box = slide.shapes.add_textbox(Inches(10.35), Inches(7.16), Inches(2.45), Inches(0.2))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Lane Binary Classifier · {page:02d}"
    set_run(run, 8, False, MUTED)


def add_card(slide, x, y, w, h, title: str | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    if title:
        box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.12), w - Inches(0.36), Inches(0.32))
        p = box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = title
        set_run(run, 13, True)
    return shape


def add_bullets(slide, items: list[str], x, y, w, h, size: int = 15) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.text = item
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(6)


def add_label(slide, text: str, x, y, w, h, color: RGBColor = BLUE, size: int = 14) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_run(run, size, True, WHITE)


def add_flow(slide, labels: list[str], x, y, w, h) -> None:
    gap = Inches(0.18)
    item_w = (w - gap * (len(labels) - 1)) / len(labels)
    for idx, label in enumerate(labels):
        left = x + idx * (item_w + gap)
        color = BLUE if idx in (0, 3, 5) else GREEN if idx in (1, 4) else ORANGE
        add_label(slide, label, left, y, item_w, h, color=color, size=11)
        if idx < len(labels) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                left + item_w - Inches(0.03),
                y + Inches(0.18),
                Inches(0.22),
                h - Inches(0.36),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LINE
            arrow.line.fill.background()


def add_code(slide, code: str, x, y, w, h) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    box.line.fill.background()
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code.strip()
    run.font.name = "Consolas"
    run.font.size = Pt(8.8)
    run.font.color.rgb = RGBColor(226, 232, 240)


def add_metric_table(slide, x, y, w, h) -> None:
    table = slide.shapes.add_table(4, 5, x, y, w, h).table
    headers = ["数据", "Acc", "Precision", "Recall", "F1"]
    rows = [
        ["TuSimple Test", "0.9665", "0.9301", "0.9732", "0.9511"],
        ["TuSimple Val", "0.9085", "0.8254", "0.9246", "0.8722"],
        ["ONNX Demo", "-", "p=0.9432 / 0.0000", "-", "OK"],
    ]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(226, 232, 240)
        p = cell.text_frame.paragraphs[0]
        p.text = header
        p.font.bold = True
        p.font.size = Pt(10)
    for r, row in enumerate(rows, 1):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(10)


def add_image_or_placeholder(slide, path: Path, x, y, w, h, label: str) -> None:
    if path.exists():
        slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    else:
        add_card(slide, x, y, w, h)
        add_label(slide, label, x + Inches(0.25), y + h / 2 - Inches(0.18), w - Inches(0.5), Inches(0.36), MUTED)


def title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_label(slide, "智能驾驶感知项目", Inches(0.78), Inches(0.7), Inches(1.7), Inches(0.38), BLUE, 12)
    box = slide.shapes.add_textbox(Inches(0.78), Inches(1.65), Inches(8.6), Inches(1.2))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "基于前视图像 ROI 的\n车道区域二分类感知模型"
    set_run(run, 34, True)
    sub = slide.shapes.add_textbox(Inches(0.82), Inches(3.25), Inches(8.2), Inches(0.7))
    run2 = sub.text_frame.paragraphs[0].add_run()
    run2.text = "Lane In/Out Binary Classification · TensorBoard · ONNX Runtime"
    set_run(run2, 17, False, MUTED)
    add_flow(
        slide,
        ["ROI", "CNN", "F1", "TensorBoard", "ONNX", "Runtime"],
        Inches(0.82),
        Inches(4.75),
        Inches(10.6),
        Inches(0.58),
    )
    add_footer(slide, 1)


def build_presentation() -> None:
    prs = Presentation()
    prs.slide_width = WIDE_WIDTH
    prs.slide_height = WIDE_HEIGHT

    title_slide(prs)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "任务背景与目标", "把车道线检测数据转化为车辆前方 ROI 的二分类判断")
    add_card(slide, Inches(0.7), Inches(1.55), Inches(5.8), Inches(4.65), "作业要求")
    add_bullets(
        slide,
        [
            "输入：车载前视图像中的局部 ROI",
            "输出：车道内 in_lane / 车道外 out_lane",
            "目标：快速验证感知环境下的基础决策能力",
            "加分：TensorBoard、ONNX、ONNX Runtime、图像增强",
        ],
        Inches(1.0),
        Inches(2.05),
        Inches(5.15),
        Inches(3.7),
    )
    add_card(slide, Inches(6.85), Inches(1.55), Inches(5.8), Inches(4.65), "项目交付")
    add_bullets(
        slide,
        [
            "可上传 GitHub 的标准工程结构",
            "训练、推理、导出和数据转换脚本",
            "实验报告 Word 与 PPT 设计文档",
            "后续可替换为 TuSimple 真实训练结果",
        ],
        Inches(7.15),
        Inches(2.05),
        Inches(5.1),
        Inches(3.7),
    )
    add_footer(slide, 2)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "技术路线总览", "从开源车道线标注到可部署的二分类模型")
    add_flow(
        slide,
        ["TuSimple\n数据集", "JSON\n标注解析", "ROI\n样本构建", "CNN\n训练", "TensorBoard\n可视化", "ONNX\n导出", "Runtime\n推理"],
        Inches(0.6),
        Inches(2.25),
        Inches(12.0),
        Inches(0.85),
    )
    add_bullets(
        slide,
        [
            "核心思路：不直接预测完整车道线，而是判断局部道路 ROI 是否处于当前车道区域内。",
            "工程闭环：数据转换、模型训练、指标可视化、模型导出、跨框架推理验证。",
            "当前状态：TuSimple 真实数据训练、测试、ONNX 导出和 Runtime 推理已完成。",
        ],
        Inches(1.0),
        Inches(4.0),
        Inches(11.0),
        Inches(1.4),
        16,
    )
    add_footer(slide, 3)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "数据集与 ROI 构建", "将车道线检测标注转换为二分类图像样本")
    add_card(slide, Inches(0.65), Inches(1.55), Inches(4.0), Inches(4.75), "TuSimple 标注")
    add_code(
        slide,
        """
{
  "raw_file": "clips/.../20.jpg",
  "h_samples": [160, 170, ...],
  "lanes": [[-2, 620, 625, ...], ...]
}
        """,
        Inches(0.95),
        Inches(2.05),
        Inches(3.4),
        Inches(1.55),
    )
    add_bullets(
        slide,
        ["raw_file：图像路径", "h_samples：固定 y 坐标", "lanes：车道线 x 坐标", "-2：无有效标注点"],
        Inches(0.95),
        Inches(3.9),
        Inches(3.4),
        Inches(1.8),
        13,
    )
    add_card(slide, Inches(4.95), Inches(1.55), Inches(7.7), Inches(4.75), "转换策略")
    add_flow(
        slide,
        ["估计\n车道中心", "中心裁剪\nin_lane", "左右偏移\nout_lane", "ImageFolder\n训练结构"],
        Inches(5.35),
        Inches(2.25),
        Inches(6.85),
        Inches(0.75),
    )
    add_bullets(
        slide,
        [
            "正样本：以当前车道中心裁剪车辆前方 ROI。",
            "负样本：在车道中心基础上向左/右偏移裁剪。",
            "输出目录：train/val 下分别包含 in_lane 与 out_lane。",
        ],
        Inches(5.35),
        Inches(3.55),
        Inches(6.85),
        Inches(1.6),
        15,
    )
    add_footer(slide, 4)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "模型结构设计", "轻量 CNN：适合作业验证和快速推理")
    add_flow(
        slide,
        ["Input\n160×96", "ConvBlock\n3→24", "ConvBlock\n24→48", "ConvBlock\n48→96", "ConvBlock\n96→128", "AvgPool\n+ Dropout", "Linear\n+ Sigmoid"],
        Inches(0.55),
        Inches(1.75),
        Inches(12.1),
        Inches(0.82),
    )
    add_card(slide, Inches(0.75), Inches(3.2), Inches(5.6), Inches(2.65), "设计理由")
    add_bullets(
        slide,
        ["卷积层提取车道线边缘、路面纹理和空间布局", "BatchNorm + ReLU 提高训练稳定性", "AdaptiveAvgPool 压缩空间特征", "单 logit 输出适配二分类任务"],
        Inches(1.05),
        Inches(3.75),
        Inches(5.0),
        Inches(1.75),
        13,
    )
    add_code(
        slide,
        """
logits = model(images)
loss = BCEWithLogitsLoss(logits, targets)
prob = sigmoid(logit)
class = in_lane if prob >= 0.5 else out_lane
        """,
        Inches(6.75),
        Inches(3.2),
        Inches(5.7),
        Inches(2.65),
    )
    add_footer(slide, 5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "训练流程与数据增强", "训练集增强，验证集保持稳定评估")
    add_card(slide, Inches(0.7), Inches(1.45), Inches(5.7), Inches(4.95), "训练配置")
    add_bullets(
        slide,
        ["输入尺寸：160 × 96", "优化器：AdamW", "学习率：1e-3", "Batch size：32", "训练轮数：20 epochs", "正式模型：best.pt，epoch 18"],
        Inches(1.05),
        Inches(1.95),
        Inches(5.0),
        Inches(3.7),
        15,
    )
    add_card(slide, Inches(6.75), Inches(1.45), Inches(5.9), Inches(4.95), "增强策略")
    add_bullets(
        slide,
        ["水平翻转：模拟左右车道位置变化", "亮度扰动：模拟不同光照条件", "对比度扰动：增强路面纹理鲁棒性", "验证集不增强：保证指标可比较"],
        Inches(7.1),
        Inches(1.95),
        Inches(5.2),
        Inches(2.2),
        15,
    )
    add_code(
        slide,
        """
image = FLIP_LEFT_RIGHT(image)
brightness = uniform(0.65, 1.35)
contrast = uniform(0.8, 1.2)
        """,
        Inches(7.1),
        Inches(4.5),
        Inches(5.15),
        Inches(1.1),
    )
    add_footer(slide, 6)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "TensorBoard 训练过程", "记录 train / val 的 loss、accuracy、precision、recall、F1")
    add_card(slide, Inches(0.65), Inches(1.45), Inches(5.85), Inches(4.9), "观察重点")
    add_bullets(
        slide,
        ["loss 是否下降", "accuracy / F1 是否上升", "train 与 val 是否明显分叉", "曲线截图可放入实验报告和答辩 PPT"],
        Inches(1.0),
        Inches(2.0),
        Inches(5.1),
        Inches(2.4),
        16,
    )
    add_code(
        slide,
        "tensorboard --logdir outputs\\tusimple_run\\tensorboard",
        Inches(1.0),
        Inches(4.85),
        Inches(5.1),
        Inches(0.75),
    )
    add_card(slide, Inches(6.85), Inches(1.45), Inches(5.8), Inches(4.9), "TuSimple 训练现象")
    add_bullets(
        slide,
        ["20 个 epoch 完成真实数据训练", "best.pt 出现在 epoch 18", "验证集 F1 = 0.8722", "测试集 F1 = 0.9511"],
        Inches(7.2),
        Inches(2.0),
        Inches(5.0),
        Inches(2.4),
        16,
    )
    add_footer(slide, 7)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "ONNX 导出与 ONNX Runtime 推理", "验证模型可脱离 PyTorch 独立运行")
    add_flow(
        slide,
        ["best.pt\nPyTorch", "export_onnx.py", "lane_binary_classifier.onnx", "predict_onnx.py", "ONNX Runtime\nCPU 推理"],
        Inches(0.8),
        Inches(1.8),
        Inches(11.7),
        Inches(0.82),
    )
    add_card(slide, Inches(0.75), Inches(3.3), Inches(5.75), Inches(2.65), "推理结果")
    add_bullets(
        slide,
        ["test in_lane：probability_in_lane = 0.9432", "test out_lane：probability_in_lane = 0.0000", "误判样例：out_lane 被判为 in_lane，p=0.7757", "provider = CPUExecutionProvider"],
        Inches(1.05),
        Inches(3.85),
        Inches(5.1),
        Inches(1.45),
        14,
    )
    add_code(
        slide,
        """
python -m lane_binary_classifier.predict_onnx ^
  --onnx-model outputs\\tusimple_run\\lane_binary_classifier.onnx ^
  --image data\\tusimple_binary\\test\\in_lane\\..._center.png
        """,
        Inches(6.85),
        Inches(3.3),
        Inches(5.75),
        Inches(2.65),
    )
    add_footer(slide, 8)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "实验结果分析", "TuSimple 独立测试集 F1 达到 0.9511")
    add_metric_table(slide, Inches(0.7), Inches(1.55), Inches(7.0), Inches(1.8))
    add_card(slide, Inches(8.15), Inches(1.55), Inches(4.45), Inches(4.75), "结果解读")
    add_bullets(
        slide,
        ["测试集共 1224 个 ROI 样本", "Accuracy = 0.9665，F1 = 0.9511", "Recall = 0.9732，车道内召回较高", "FP=30 大于 FN=11，后续可优化困难负样本"],
        Inches(8.45),
        Inches(2.05),
        Inches(3.85),
        Inches(2.5),
        14,
    )
    add_card(slide, Inches(0.7), Inches(3.85), Inches(7.0), Inches(2.45), "混淆矩阵")
    add_bullets(
        slide,
        ["TP=399：真实 in_lane，预测 in_lane", "TN=784：真实 out_lane，预测 out_lane", "FP=30：真实 out_lane，误判 in_lane", "FN=11：真实 in_lane，误判 out_lane"],
        Inches(1.05),
        Inches(4.35),
        Inches(6.3),
        Inches(1.4),
        14,
    )
    add_footer(slide, 9)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "问题记录与解决方案", "工程问题也是项目完成度的一部分")
    add_card(slide, Inches(0.7), Inches(1.45), Inches(12.0), Inches(4.9))
    rows = [
        ("pip 命令误连", "分开执行 pip install -r 和 pip install -e"),
        ("TensorBoard 提示缺少 TensorFlow", "PyTorch 曲线仍可正常显示，无需安装 TensorFlow"),
        ("ONNX 缺少依赖", "补充 onnx、onnxscript、onnxruntime"),
        ("Windows 中文终端导出编码问题", "使用 legacy ONNX export，dynamo=False"),
        ("GitHub workflow 权限不足", "移除非必要 CI 文件后成功 push"),
    ]
    y = Inches(1.85)
    for idx, (problem, solution) in enumerate(rows, 1):
        add_label(slide, str(idx), Inches(1.05), y, Inches(0.35), Inches(0.35), BLUE, 10)
        add_bullets(slide, [f"{problem} → {solution}"], Inches(1.55), y - Inches(0.03), Inches(10.5), Inches(0.42), 13)
        y += Inches(0.75)
    add_footer(slide, 10)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "方案优缺点", "轻量二分类方案适合快速验证，也有明确边界")
    add_card(slide, Inches(0.7), Inches(1.45), Inches(5.85), Inches(4.95), "优点")
    add_bullets(
        slide,
        ["模型轻量，训练和推理速度快", "ROI 聚焦车辆前方关键区域", "支持 TensorBoard 可视化", "支持 ONNX 部署验证", "数据转换逻辑清晰，易复现"],
        Inches(1.05),
        Inches(1.95),
        Inches(5.1),
        Inches(3.4),
        14,
    )
    add_card(slide, Inches(6.85), Inches(1.45), Inches(5.85), Inches(4.95), "不足")
    add_bullets(
        slide,
        ["无法输出精确车道线位置", "ROI 标签质量依赖裁剪策略", "部分困难 out_lane 会被误判", "未利用连续帧时序信息", "复杂天气、遮挡、弯道仍需进一步验证"],
        Inches(7.2),
        Inches(1.95),
        Inches(5.1),
        Inches(3.4),
        14,
    )
    add_footer(slide, 11)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "总结与下一步工作", "真实数据训练与部署验证均已完成")
    add_card(slide, Inches(0.7), Inches(1.45), Inches(5.85), Inches(4.95), "当前完成")
    add_bullets(
        slide,
        ["完成 TuSimple 数据转换与训练", "测试集 Accuracy=0.9665，F1=0.9511", "完成 TensorBoard 训练记录", "完成 ONNX Runtime 独立推理", "完成 Word/PPT 作业材料"],
        Inches(1.05),
        Inches(1.95),
        Inches(5.1),
        Inches(3.4),
        14,
    )
    add_card(slide, Inches(6.85), Inches(1.45), Inches(5.85), Inches(4.95), "下一阶段")
    add_bullets(
        slide,
        ["补充更多误差样例截图", "尝试 MobileNet / ResNet18", "调整阈值降低 FP", "统计 ONNX Runtime 推理耗时", "加入更多天气/光照增强"],
        Inches(7.2),
        Inches(1.95),
        Inches(5.1),
        Inches(3.4),
        14,
    )
    add_footer(slide, 12)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build_presentation()
